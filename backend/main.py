import io
import os
import json
from typing import Dict, Tuple
from fastapi import FastAPI, HTTPException, Query, Header
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, validator
from typing import List, Dict, Union, Optional
from pptx import Presentation
from signup import generate_next_id, create_user_folders, save_users
from ppt import create_slide

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Adjust this to your React app's URL
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ContentValueType = Union[str, List[str], List[Dict[str, Union[str, int]]]]
ContentValueType = Union[str, List[str], List[List[Union[str, int]]], List[Dict[str, Union[str, int]]]]

class ContentItem(BaseModel):
    type: str
    value: Optional[ContentValueType] = None
    data: Optional[ContentValueType] = None
    headers: Optional[List[str]] = None
    graphType: str = None

class SlideData(BaseModel):
    title: str
    content: List[ContentItem]

class PresentationData(BaseModel):
    slides: List[SlideData]

class UserLogin(BaseModel):
    email: str
    password: str

def ensure_directory(path):
    os.makedirs(path, exist_ok=True)

# Function 1: Save data to a JSON file
def save_to_json(path: str, filename: str, data: dict):
    full_path = os.path.join(path, filename)
    try:
        os.makedirs(os.path.dirname(full_path), exist_ok=True)
        with open(full_path, 'w') as file:
            json.dump(data, file, indent=2)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error saving file: {str(e)}")

# Function 2: Read data from a JSON file
def read_from_json(path: str, filename: str) -> Tuple[Dict, int, str]:
    full_path = os.path.join(path, filename)
    try:
        with open(full_path, 'r') as file:
            return json.load(file), 200, ""
    except FileNotFoundError:
        return {}, 404, "File not found"
    except json.JSONDecodeError:
        return {}, 400, "Invalid JSON file"
    except Exception as e:
        return {}, 500, f"Error reading file: {str(e)}"

def filter_slides(slides, selected_cards):
    if not selected_cards:
        return {"slides": slides}

    filtered_slides = []

    for slide in slides:
        filtered_content = []
        for content in slide["content"]:
            filtered_data = [
                row for row in content["data"]
                if row[0] in selected_cards
            ]

            if filtered_data:
                filtered_content.append({
                    "headers": content["headers"],
                    "data": filtered_data,
                    "type": content["type"]
                })

        if filtered_content:
            filtered_slides.append({
                "title": slide["title"],
                "content": filtered_content
            })

    return {"slides": filtered_slides}

def user_exists(user_id):
    users, status_code, error_message = read_from_json("./jsons/users/", "users.json")
    if status_code != 200:
        if status_code == 404:
            raise HTTPException(status_code=500, detail="Users file not found")
        elif status_code == 400:
            raise HTTPException(status_code=500, detail="Invalid users file")
        else:
            raise HTTPException(status_code=500, detail=error_message)

    for user in users:
        if user.get('id') == user_id:
            return {"message": "Login successful", "data": user_id}

    return None

def convert_ai_use_json_list(input_list):
    converted_list = []
    
    for input_json in input_list:
        # Extract risks dictionary for easier access
        risks = input_json.get('risks', {})

        converted_json = {
            "name": input_json.get('value_area_name'),
            "applicationTypes": input_json.get('category_name'),
            "horizon": input_json.get('ai_use_class')+ " " +input_json.get('horizonClassification'),
            "infoExposure": risks.get('infoExposure'),
            "dataSensitivity": risks.get('infoSensitivity'),
            "complexityOfUsage": input_json.get('riskComplexity'),
            "riskRating": risks.get('riskRating'),
            "description":input_json.get('description'),
            "reasoning": [risks.get('reasoning')] if risks.get('reasoning') else [],
            "level_of_specialisation_interactionact":risks.get('level_of_specialisation_interaction'),
            "app_operation":risks.get('app_operation'),
        }
        
        converted_list.append(converted_json)
    return converted_list

def extract_applications(data):
    applications = []
    id_counter = 1  # Initialize the counter

    for key, value in data.items():  # Changed from item() to items() to fix syntax
        for application in value['applications']:
            application['id'] = id_counter  # Add the incremental id
            applications.append(application)  # Use append instead of extend
            id_counter += 1  # Increment the counter for the next application

    return applications

def transform_tech_reasoning_json(user_id):

    input_data, input_code, input_message = read_from_json(f"./jsons/{user_id}/output/", "tech_quad_dict.json")
    question_data, question_code, question_message = read_from_json(f"./jsons/{user_id}/output/", "tech_quad_question.json")

    if input_code != 200:
        raise HTTPException(status_code=input_code, detail=input_message)
    if question_code != 200:
        raise HTTPException(status_code=question_code, detail=question_message)

    if input_data is None:
        return

    transformed_data = []
    
    for _, item in input_data.items():
        transformed_item = {
            "title": f"{item['new_app_name']} - {item['category_name']}",
            "description": item['description'] if 'description' in item else '',
            "defaultInfo": [
                {
                    "icon": "🧹",
                    "title": "Linked AI Value",
                    "options": [
                        {"value": item['value_area_name'], "selected": True}
                    ]
                },
                {
                    "icon": "📊",
                    "title": "AI Responsiible Use",
                    "options": [
                        {"value": "Everyday use", "selected": item['ai_use_class'] == "Everyday use"},
                        {"value": "Augmentation", "selected": item['ai_use_class'] == "Augmentation"},
                        {"value": "Transformation", "selected": item['ai_use_class'] == "Transformation"}
                    ]
                },
                {
                    "icon": "🔒",
                    "title": "Quadrant Mapping",
                    "options": [
                        {"value": "Quadrant-1", "selected": item['ai_enablement_quadrant'] == "Quadrant 1"},
                        {"value": "Quadrant-2", "selected": item['ai_enablement_quadrant'] == "Quadrant 2"},
                        {"value": "Quadrant-3", "selected": item['ai_enablement_quadrant'] == "Quadrant 3"},
                        {"value": "Quadrant-4", "selected": item['ai_enablement_quadrant'] == "Quadrant 4"}
                    ]
                }
            ],
            "quadrant_reasoning": item['quadrant_reasoning'] if 'quadrant_reasoning' in item else '',
            "questions":question_data
        }
        transformed_data.append(transformed_item)
    
    return  transformed_data

def create_the_ppt(prs, data = None):
    prs = Presentation()
    # ======================================== slide_1 ================================================== table
    app_list = list(data.items())
    for i in range(0, len(app_list), 3):
        
        others = prs.slide_masters[2].slide_layouts[4] 
        slide = prs.slides.add_slide(others)
        subheading = slide.placeholders[13]
        subheading.text = ''
        
        title = slide.shapes.title
        title.text = 'AI App Names and Descriptions'
        
        # Remove all shapes except for the title and subtitle
        for shape in slide.shapes:
            if not shape.has_text_frame or shape == slide.shapes.title or shape == slide.placeholders[13]:
                continue
            sp = shape._element
            sp.getparent().remove(sp)
        
        rows = min(3, len(app_list) - i) + 1  # Number of apps in this table + header row
        cols = 2
        left = Inches(0.755)
        top = Inches(2.3)
        width = Inches(17.5)
        height = Inches(1.2)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Set column widths
        table.columns[0].width = Inches(3.75)
        table.columns[1].width = Inches(13.25)
        
        # Set table headers
        table.cell(0, 0).text = 'App Name'
        table.cell(0, 1).text = 'Description'
        
        # Add data to the table
        for j in range(min(3, len(app_list) - i)):
            app_id, app_data = app_list[i + j]
            table.cell(j + 1, 0).text = app_data['new_app_name']
            table.cell(j + 1, 1).text = app_data['description']
            
    # ======================================== slide_2 ================================================== 
    # Add slides with tables for each application (up to 5 rows per table)
    app_list = list(data.items())
    for i in range(0, len(app_list), 5):
        others = prs.slide_masters[2].slide_layouts[4] 
        slide = prs.slides.add_slide(others)
        title = slide.shapes.title
        title.text = "AI Applications Summary"
        subheading = slide.placeholders[13]
        subheading.text = 'Short summary of applications discussed and finalised throughout the tool'
        
        # Remove all shapes except for the title and subtitle
        for shape in slide.shapes:
            if not shape.has_text_frame or shape == slide.shapes.title or shape == slide.placeholders[13]:
                continue
            sp = shape._element
            sp.getparent().remove(sp)
        
        rows = min(5, len(app_list) - i) * 2 + 1  # Number of apps in this table * 2 (for Class and Reasoning) + header row  # Number of apps in this table + header row
        cols = 5
        left = Inches(0.775)
        top = Inches(1.9)
        width = Inches(17.5)
        height = Inches(1.2)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Set column widths
        table.columns[0].width = Inches(1.5)
        table.columns[1].width = Inches(3.55)
        table.columns[2].width = Inches(3.55)
        table.columns[3].width = Inches(3.55)
        table.columns[4].width = Inches(3.55)
        
        # Set table headers
        table.cell(0, 0).text = 'Sr. No.'
        table.cell(0, 1).text = 'App Name'
        table.cell(0, 2).text = 'Linked AI Value'
        table.cell(0, 3).text = 'AI Responsible Use'
        table.cell(0, 4).text = 'AI Tech Enablement'
        
        # Add data to the table
        for j in range(min(5, len(app_list) - i)):
            app_id, app_data = app_list[i + j]
            # Class row
            table.cell(j * 2 + 1, 0).text = str(app_id)
            table.cell(j * 2 + 1, 1).text = app_data['new_app_name']
            table.cell(j * 2 + 1, 2).text = app_data['value_area_name']
            table.cell(j * 2 + 1, 3).text = app_data['ai_use_class']
            table.cell(j * 2 + 1, 4).text = app_data['enablement_route']
            # Reasoning row
            table.cell(j * 2 + 2, 2).text = "" #app_data['description']
            table.cell(j * 2 + 2, 3).text = app_data['category_name']
            table.cell(j * 2 + 2, 4).text = app_data['ai_enablement_quadrant']

    # ======================================== slide_2 ==================================================
    divider_slide_layout = prs.slide_masters[0].slide_layouts[2] 
    slide = prs.slides.add_slide(divider_slide_layout)
    title = slide.placeholders[14]
    title.text = "Thank You"
    
    title_text_frame = title.text_frame
    title_text_frame.text = "Thank You"
    title_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    return prs

# Pydantic model for request body
class DataModel(BaseModel):
    data: dict

class UserLoginData(BaseModel):
    email: str
    password: str
class UserLogin(BaseModel):
    data: UserLoginData

class DataModelList(BaseModel):
    data: list

class GeneratePPTRequest(BaseModel):
    type: str  
    selected_cards: List[dict] 

@app.post("/save-business-context")
async def save_business_context(data: DataModel, User_id: str = Header(...)):
    user_check_result = user_exists(User_id)
    if user_check_result is None:
        raise HTTPException(status_code=403, detail="User not found")

    save_to_json(f"./jsons/{User_id}/input/", "business_context.json", data.data)

    suggestion_output, status_code, error_message = read_from_json(f"./jsons/{User_id}/output/", "suggestion_output.json")
    if status_code != 200:
        raise HTTPException(status_code=status_code, detail=error_message)
    return suggestion_output

@app.post("/save-suggestions")
async def save_suggestions(data: DataModel, User_id: str = Header(...)):
    user_check_result = user_exists(User_id)
    if user_check_result is None:
        raise HTTPException(status_code=403, detail="User not found")
    save_to_json(f"./jsons/{User_id}/input/", "suggestion_input.json", data.data)

    opportunities, opp_status, opp_error = read_from_json(f"./jsons/{User_id}/output/", "opportunities.json")

    if opp_status != 200:
        raise HTTPException(status_code=opp_status, detail=opp_error)

    return opportunities

@app.get("/fetch-opportunities")
async def fetch_opportunities(User_id: str = Header(...)):
    user_check_result = user_exists(User_id)
    if user_check_result is None:
        raise HTTPException(status_code=403, detail="User not found")

    opportunities, opp_status, opp_error = read_from_json(f"./jsons/{User_id}/output/", "opportunities.json")

    if opp_status != 200:
        raise HTTPException(status_code=opp_status, detail=opp_error)

    return opportunities
    
@app.get("/fetch-ai-use")
async def fetch_ai_use(User_id: str = Header(...)):
    user_check_result = user_exists(User_id)
    if user_check_result is None:
        raise HTTPException(status_code=403, detail="User not found")

    ai_use, ai_status, ai_error = read_from_json(f"./jsons/{User_id}/output/","ai_use_gpt_output.json")
    result = convert_ai_use_json_list(ai_use["applications"])
    return result

@app.get("/fetch-ai-strategy")
async def fetch_ai_strategy(User_id: str = Header(...)):
    user_check_result = user_exists(User_id)
    if user_check_result is None:
        raise HTTPException(status_code=403, detail="User not found")

    ai_strategy, ai_status, ai_error = read_from_json(f"./jsons/{User_id}/output/", "ai_strategy_output.json")
    application_list = extract_applications(ai_strategy)
    return application_list

@app.get("/fetch-tech-reasoning")
async def fetch_tech_reasoning(User_id: str = Header(...)):
    user_check_result = user_exists(User_id)
    if user_check_result is None:
        raise HTTPException(status_code=403, detail="User not found")
    transformed_data = transform_tech_reasoning_json(User_id)
    save_to_json(f"./jsons/{User_id}/output/", "tech_enablement_values.json", transformed_data)

    return transformed_data

@app.post("/save-ai-strategy")
async def save_ai_strategy(data: DataModelList, User_id: str = Header(...)):
    user_check_result = user_exists(User_id)
    if user_check_result is None:
        raise HTTPException(status_code=403, detail="User not found")

    save_to_json(f"./jsons/{User_id}/input/", "ai_strategy_input.json", data.data)
    ai_use, ai_status, ai_error = read_from_json(f"./jsons/{User_id}/output/", "ai_use.json")
    if ai_status != 200:
        raise HTTPException(status_code=ai_status, detail=error_message)
    return ai_use

@app.post("/save-tech-reasoning")
async def save_tech_reasoning(data: DataModel, User_id: str = Header(...)):
    user_check_result = user_exists(User_id)
    if user_check_result is None:
        raise HTTPException(status_code=403, detail="User not found")

    save_to_json(f"./jsons/{User_id}/input/", "tech_reasoning.json", data.data)
    final_result, status_code, error_message = read_from_json(f"./jsons/{User_id}/output/", "result.json")
    if status_code != 200:
        raise HTTPException(status_code=status_code, detail=error_message)
    return final_result

@app.post("/generate-ppt")
async def generate_ppt(request: GeneratePPTRequest, User_id: str = Header(...)):
    user_check_result = user_exists(User_id)
    if user_check_result is None:
        raise HTTPException(status_code=403, detail="User not found")

    if request.type.lower() not in ["short", "comprehensive"]:
        raise HTTPException(status_code=400, detail="Invalid type. Use 'short' or 'comprehensive'.")
    
    json_file = "short_table.json" if request.type == "short" else "table.json"
    filename = "short_presentation.pptx" if request.type == "short" else "presentation.pptx"

    save_to_json(f"./jsons/{User_id}/input/", "final_selected_result.json", request.selected_cards)
    # json_data, status_code, error_message = read_from_json(f"./jsons/{User_id}/output/", json_file)
    json_data, status_code, error_message = read_from_json(f"./jsons/{User_id}/gpt_output/", "final_result.json")

    if status_code != 200:
        raise HTTPException(status_code=status_code, detail=error_message)

    prs = create_the_ppt(prs, json_data)

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    return StreamingResponse(
        ppt_io,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

@app.post("/login")
def login(data: UserLogin):
    users, status_code, error_message = read_from_json("./jsons/users/", "users.json")
    if status_code != 200:
        raise HTTPException(status_code=status_code, detail=error_message)
    for u in users:
        if u['email'] == data.data.email and u['password'] == data.data.password:
            user_details = u.copy()
            del user_details['password']
            return {"message": "Login successful", "data": user_details}
    raise HTTPException(status_code=400, detail="Invalid email or password")

@app.post("/signup")
async def signup(data: DataModel):
    users, status_code, error_message = read_from_json("./jsons/users/", "users.json")
    
    if any(existing_user['email'] == data.data["email"] for existing_user in users):
        raise HTTPException(status_code=400, detail="User with this email already exists")
    
    if any(existing_user['name'] == data.data["name"] for existing_user in users):
        raise HTTPException(status_code=400, detail="Username already taken")
    
    existing_ids = [u['id'] for u in users]
    new_user_id = generate_next_id(existing_ids)
    
    new_user = {
        "id": new_user_id,
        "name": data.data["name"],
        "email": data.data["email"],
        "password": data.data["password"]  # Note: In production, hash the password!
    }
    
    users.append(new_user)
    # save_users(users)
    save_to_json(f"./jsons/users/", "users.json", users)
    
    create_user_folders(new_user_id)
    
    return {"message": "User created successfully", "data": {"id":new_user["id"], "name":new_user["name"], "email":new_user["email"]}}

@app.get("/")
async def read_root():
    return {"message": "Welcome to the FastAPI application!"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)