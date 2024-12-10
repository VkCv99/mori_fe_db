from pydantic import BaseModel, validator
from typing import List, Dict, Union, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor

# ContentValueType = Union[str, List[str], List[Dict[str, Union[str, int]]]]
ContentValueType = Union[str, List[str], List[List[Union[str, int]]], List[Dict[str, Union[str, int]]]]

class ContentItem(BaseModel):
    type: str
    value: Optional[ContentValueType] = None
    data: Optional[ContentValueType] = None
    headers: Optional[List[str]] = None
    graphType: str = None

class SlideData(BaseModel):
    title: str
    content: List[ContentItem]

class PresentationData(BaseModel):
    slides: List[SlideData]

def text_fits(text_frame, max_size):
    """
    Check if the text fits within the text frame.
    This is an approximation and may not be 100% accurate in all cases.
    """
    if not text_frame.text:
        return True

    text_frame.fit_text(font_family='Arial', max_size=max_size, bold=False, italic=False)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size < 6:  # If font size is too small, consider it doesn't fit
                return False
    return True

def fit_text(text_frame, text_type='normal'):
    """
    Fit text into the text frame by adjusting size and applying wrapping.
    
    :param text_frame: The text frame to fit text into
    :param text_type: Type of text ('title', 'subtitle', or 'normal')
    """
    if text_type == 'title':
        max_size, min_size = 24, 20
    elif text_type == 'subtitle':
        max_size, min_size = 14, 12
    else:  # normal text
        max_size, min_size = 11, 9

    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # Start with the maximum size
    size = max_size
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(size)

    while not text_fits(text_frame, size) and size > min_size:
        size = max(size * 0.9, min_size)  # Reduce by 10% each iteration
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = size

    # If still doesn't fit, try to wrap (for titles and subtitles)
    if not text_fits(text_frame, size) and text_type in ['title', 'subtitle']:
        words = text_frame.text.split()
        mid = len(words) // 2
        text_frame.text = ' '.join(words[:mid]) + '\n' + ' '.join(words[mid:])

        # Readjust size after wrapping
        while not text_fits(text_frame, size) and size > min_size:
            size = max(size * 0.9, min_size)
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = size

    # For normal text, if it still doesn't fit, truncate with ellipsis
    if not text_fits(text_frame, size) and text_type == 'normal':
        while not text_fits(text_frame, size) and len(text_frame.text) > 3:
            text_frame.text = text_frame.text[:-4] + '...'

    return Pt(size)

def add_chart(slide, data, chart_type):
    chart_data = CategoryChartData()
    chart_data.categories = [item['name'] for item in data]
    chart_data.add_series('Series 1', [item['value'] for item in data])

    # Adjust chart size and position
    x, y = Inches(2.5), Inches(3)
    cx, cy = Inches(5), Inches(2.5)  # Reduced size
    chart = slide.shapes.add_chart(
        chart_type, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(8)  # Adjust legend font size

    # Adjust font sizes and other chart properties
    if chart.chart_title:
        chart.chart_title.text_frame.text = "Chart Title"  # You can customize this
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

    # Adjust series data labels
    for series in chart.series:
        if hasattr(series, 'data_labels'):
            series.data_labels.font.size = Pt(8)
            if chart_type != XL_CHART_TYPE.PIE:
                series.data_labels.position = XL_DATA_LABEL_POSITION.CENTER

    # Adjust axes based on chart type
    if chart_type != XL_CHART_TYPE.PIE:
        # For charts with category and value axes
        if hasattr(chart, 'category_axis') and hasattr(chart.category_axis, 'tick_labels'):
            chart.category_axis.tick_labels.font.size = Pt(8)

        if hasattr(chart, 'value_axis') and hasattr(chart.value_axis, 'tick_labels'):
            chart.value_axis.tick_labels.font.size = Pt(8)

    else:
        # For pie charts
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(8)
        data_labels.position = XL_DATA_LABEL_POSITION.CENTER

    # Additional formatting
    plot = chart.plots[0]
    plot.has_data_labels = True
    if hasattr(plot, 'data_labels'):
        plot.data_labels.font.size = Pt(8)
        plot.data_labels.font.color.rgb = RGBColor(0, 0, 0)  # Black color
        
def create_slide(prs:Presentation, slide_data:SlideData):
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # Title
    title = slide.shapes.title
    title.text = slide_data.title
    fit_text(title.text_frame, 'title')

    content_top = Inches(1.5)
    for item in slide_data.content:
        if item.type == 'text':
            left, width, height = Inches(0.5), Inches(9), Inches(0.5)
            tf = slide.shapes.add_textbox(left, content_top, width, height).text_frame
            tf.text = item.value
            fit_text(tf, 'normal')
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            content_top += height + Inches(0.1)

        elif item.type == 'bullet':
            shapes = slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = item.value[0]
            for bullet in item.value[1:]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 1
            fit_text(tf, 'normal')
            content_top = Inches(5.5)  # Move to bottom half of slide

        elif item.type == 'graph':
            if item.graphType == 'bar':
                add_chart(slide, item.data, XL_CHART_TYPE.COLUMN_CLUSTERED)
            elif item.graphType == 'pie':
                add_chart(slide, item.data, XL_CHART_TYPE.PIE)
            content_top = Inches(15)  # Move to bottom half of slide
        elif item.type == 'table':
            headers = item.headers
            data = item.data
            rows, cols = len(data) + 1, len(headers)  # +1 for header row
            left, top, width, height = Inches(0.5), content_top, Inches(9), Inches(0.5 * rows)
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Populate the header row
            for col, header in enumerate(headers):
                cell = table.cell(0, col)
                cell.text = str(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(10)

            # Populate the data rows
            for row, row_data in enumerate(data, start=1):
                for col, cell_data in enumerate(row_data):
                    cell = table.cell(row, col)
                    cell.text = str(cell_data)
                    cell.text_frame.paragraphs[0].font.size = Pt(10)


            content_top += height + Inches(0.1)